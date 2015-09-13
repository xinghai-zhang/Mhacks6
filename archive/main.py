from newspaper import Article
from pptx import Presentation
from pptx.util import Inches, Pt
import re
from nltk.corpus import stopwords     
from reduction import *
import urllib



url = u'http://www.nytimes.com/2015/09/13/business/at-wework-an-idealistic-startup-clashes-with-its-cleaners.html?ref=business'
article = Article(url)
article.download()
article.html
article.parse()	
#reducing sentences
reduction = Reduction()
text = article.text
reduction_ratio = 0.05
reduced_text = reduction.reduce(text, reduction_ratio)	
#start pptx
prs = Presentation()
SLD_LAYOUT_TITLE_AND_CONTENT = 1	
#title page
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)			#top image
urllib.urlretrieve(str(article.top_image), "1.png")
img_path = "1.png"
height = Inches(2)
left = Inches(3)
top = Inches(5)
pic = slide.shapes.add_picture(img_path, left, top, height=height)	
#add title
title = slide.shapes.title
subtitle = slide.placeholders[1]
author_name = 'Xinghai Zhang'
title.text = article.title
subtitle.text = author_name	
#subpage
for sentence in reduced_text:
	if "Advertisement" not in sentence: 
		slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
		slide = prs.slides.add_slide(slide_layout)
		shapes = slide.shapes		
		body_shape = shapes.placeholders[1]			
		tf = body_shape.text_frame
		tf.text = sentence	
#end page
end_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(end_slide_layout)
left = width = Inches(4)
txBox = slide.shapes.add_textbox(left, Inches(2), width, Inches(2.5))
tf = txBox.text_frame			
p = tf.add_paragraph()
p.text = "The End"
p.font.bold = True
p.font.size = Pt(50)	
p = tf.add_paragraph()
p.text = "  Thank You"
p.font.size = Pt(30)			
prs.save('test.pptx')

 
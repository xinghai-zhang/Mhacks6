"""
Routes and views for the flask application.
"""

from datetime import datetime
from flask import render_template, request, Flask, redirect, url_for
from FlaskWebProject1 import app
import os
import requests
from newspaper import Article
from pptx import Presentation
from pptx.util import Inches, Pt    
from reduction import *
import urllib
from FlaskWebProject1 import APP_STATIC




@app.route('/')
@app.route('/home')
def home():
    """Renders the home page."""
    return render_template(
        'index.html',
        title='Home Page',
        year=datetime.now().year,
    )

@app.route('/contact')
def contact():
    """Renders the contact page."""
    return render_template(
        'contact.html',
        title='Contact',
        year=datetime.now().year,
        message='Your contact page.'
    )

@app.route('/about')
def about():
    """Renders the about page."""
    return render_template(
        'about.html',
        title='About',
        year=datetime.now().year,
        message='Your application description page.'
    )

@app.route('/', methods=['GET', 'POST'])
def index():
    errors = []
    results = {}
    if request.method == "POST":
        # get url that the person has entered
        try:
            url = request.form['url']
            r = requests.get(url)
        except:
            errors.append(
                "Unable to get URL. Please make sure it's valid and try again."
            )
            return render_template('index.html', errors=errors)
        if r:
            article = Article(url)
            article.download()
            article.parse()    
            #reducing sentences
            reduction = Reduction()
            text = article.text
            reduction_ratio = 0.05
            reduced_text = reduction.reduce(text, reduction_ratio)    
            #start pptx
            prs = Presentation()
            SLD_LAYOUT_TITLE_AND_CONTENT = 1    
            #title page
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)            

            #top image
            urllib.urlretrieve(str(article.top_image), "1.png")
            img_path = "1.png"
            height = Inches(2)
            left = Inches(3)
            top = Inches(5)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)    
            #add title
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            author_name = 'Xinghai Zhang'
            title.text = article.title
            subtitle.text = author_name    
            #subpage
            for sentence in reduced_text:
                if "Advertisement" not in sentence: 
                    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
                    slide = prs.slides.add_slide(slide_layout)
                    shapes = slide.shapes        
                    body_shape = shapes.placeholders[1]            
                    tf = body_shape.text_frame
                    tf.text = sentence    
            #end page
            end_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(end_slide_layout)
            left = width = Inches(4)
            txBox = slide.shapes.add_textbox(left, Inches(2), width, Inches(2.5))
            tf = txBox.text_frame            

            p = tf.add_paragraph()
            p.text = "The End"
            p.font.bold = True
            p.font.size = Pt(50)    
            p = tf.add_paragraph()
            p.text = "  Thank You"
            p.font.size = Pt(30)            

            prs.save(os.path.join(APP_STATIC, 'test.pptx')) 
            return redirect(url_for(APP_STATIC, filename='test.pptx'))      


    return render_template('index.html', errors=errors, results=results)

